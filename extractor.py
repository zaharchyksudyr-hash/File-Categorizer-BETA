# Імпортуються бібліотеки для роботи з різними форматами файлів:

# Path — робота зі шляхами до файлів;
# PdfReader — читання PDF;
# Document — читання DOCX;
# Presentation — читання PPTX;
# load_workbook — читання XLSX.

from __future__ import annotations
from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

# Це головна функція модуля.
# Вона:
    # визначає розширення файлу;
    # викликає відповідний метод витягування тексту.

def extract_text(path: Path, *, max_chars: int = 2_000) -> str:

    suffix = path.suffix.lower()

    if suffix == ".docx":
        return _extract_docx(path, max_chars=max_chars)
    if suffix == ".txt":
        return _extract_txt(path, max_chars=max_chars)
    if suffix == ".pdf":
        return _extract_pdf(path, max_chars=max_chars)
    if suffix == ".md":
        return _extract_txt(path, max_chars=max_chars)
    if suffix == ".xlsx":
        return _extract_xlsx(path, max_chars=max_chars)
    if suffix == ".pptx":
        return _extract_pptx(path, max_chars=max_chars)
    return ""


# Функція намагається відкрити файл у різних кодуваннях:
    # utf-8
    # utf-8-sig
    # cp1251
    # latin-1
# Якщо одне з них підходить — текст читається успішно.
# Якщо жодне не спрацювало — робиться остання спроба.
# Після цього текст обрізається до max_chars.

def _extract_txt(path: Path, *, max_chars: int) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            text = path.read_text(encoding=enc, errors="strict")
            return text[:max_chars]
        except UnicodeDecodeError:
            continue
        except OSError:
            return ""

    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return text[:max_chars]
    except OSError:
        return ""

# Документ відкривається через python-docx.
# Далі:
    # проходяться всі абзаци;
    # порожні пропускаються;
    # текст кожного абзацу додається у список.
# Наприкінці:
    # "\n".join(parts)
    # усі абзаци об'єднуються в один текст.
# Після цього текст обрізається до максимальної довжини.

def _extract_docx(path: Path, *, max_chars: int) -> str:
    try:
        doc = Document(str(path))
        parts: list[str] = []
        for para in doc.paragraphs:
                if para.text:
                    parts.append(para.text)
        text = "\n".join(parts)
        return text[:max_chars]
    except Exception:
            return ""

# Документ відкривається через PdfReader.
# Далі:
    # читаються сторінки;
    # текст витягується методом
    # page.extract_text()
# Щоб не обробляти великі документи, встановлено обмеження: max_pages = 3
# тобто аналізуються лише перші три сторінки.
# Також контролюється загальна кількість символів.
# Як тільки досягнуто max_chars — читання припиняється.

def _extract_pdf(path: Path, *, max_chars: int) -> str:
    try:
        reader = PdfReader(str(path))
        parts: list[str] = []
        total_len = 0
        max_pages = 3
        for page_index, page in enumerate(reader.pages):
            if page_index >= max_pages:
                break
            page_text = page.extract_text() or ""
            page_text = page_text.strip()
            if not page_text:
                continue
            remaining = max_chars - total_len
            if remaining <= 0:
                break
            chunk = page_text[:remaining]
            parts.append(chunk)
            total_len += len(chunk)
            if total_len >= max_chars:
                break
        return "\n".join(parts)
    except Exception:
        return ""

# Презентація відкривається через Presentation.
# Далі:
    # проходяться слайди;
    # проходяться всі об'єкти на слайді;
    # якщо об'єкт містить текст — він додається до результату.
# Для оптимізації встановлено обмеження:
    # max_slides = 10
# тобто аналізуються лише перші десять слайдів.
# Також контролюється загальна кількість символів.

def _extract_pptx(path: Path, *, max_chars: int) -> str:
    try:
        prs = Presentation(str(path))
        parts: list[str] = []
        total_len = 0
        max_slides = 10
        for slide_index, slide in enumerate(prs.slides):
            if slide_index >= max_slides:
                break
            for shape in slide.shapes:
                text = ""
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                if not text:
                    continue
                remaining = max_chars - total_len
                if remaining <= 0:
                    return "\n".join(parts)
                chunk = text[:remaining]
                parts.append(chunk)
                total_len += len(chunk)
                if total_len >= max_chars:
                    return "\n".join(parts)
        return "\n".join(parts)
    except Exception:
        return ""


# Excel-файл відкривається у режимі лише для читання:
#     read_only=True
# Це дозволяє працювати швидше та економити пам'ять.
# Далі:
#     проходяться аркуші;
#     проходяться рядки;
#     з кожного рядка вибираються непорожні клітинки;
#     клітинки об'єднуються через символ
# Щоб уникнути надто великої обробки, встановлено обмеження:
#     max_sheets = 2
#     max_rows_per_sheet = 50
# Після завершення книга закривається.

def _extract_xlsx(path: Path, *, max_chars: int) -> str:
    try:
        wb = load_workbook(filename=str(path), read_only=True, data_only=True)
        parts: list[str] = []
        total_len = 0
        max_sheets = 2
        max_rows_per_sheet = 50
        for sheet_index, sheet in enumerate(wb.worksheets):
            if sheet_index >= max_sheets:
                break
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count >= max_rows_per_sheet:
                    break
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if not values:
                    row_count += 1
                    continue
                line = " | ".join(values)
                remaining = max_chars - total_len
                if remaining <= 0:
                    wb.close()
                    return "\n".join(parts)
                chunk = line[:remaining]
                parts.append(chunk)
                total_len += len(chunk)
                if total_len >= max_chars:
                    wb.close()
                    return "\n".join(parts)
                row_count += 1

        wb.close()
        return "\n".join(parts)
    except Exception:
        return ""
